from presentation import Presentation
from apis import replace_image_with_table, replace_image, merge_cells
from shapes import Picture
from utils import Config
from pptx import Presentation as PPTXPre
from copy import deepcopy
import pdb

config = Config('/root/pptagent/pptagent-v4/template')
presentation = Presentation.from_file("/root/pptagent/pptagent-v4/template/source.pptx", config)
slide_page = presentation.slides[0]
picture = next((s for s in slide_page.shapes if isinstance(s, Picture)), None)
if picture is None:
    print("No picture found in the slide.")
else:
    shape_idx = picture.shape_idx
    table_data = [['Camera Control', 'Camera Control'],
                  ['',''], 
                  ['Pro camera system', 'Advanced dual-camera system'],
                  ['Our most advanced 48MP Fusion camera', 'Advanced 48MP Fusion camera'],
                  ['5x Telephoto camera', '2x Telephoto'],
                  ['48MP Ultra Wide camera', '12MP Ultra Wide camera'],
                  ['Up to 33 hours video playback', 'Up to 27 hours video playback']]
    # table_data = [['Use-case', '(%)'],
    #               ['---', '---'],
    #               ['Generation', '45.6%'],
    #               ['Open QA', '12.4%'],
    #               ['Brainstorming', '11.2%'],
    #               ['Chat', '8.4%'],
    #               ['Rewrite', '6.6%'],
    #               ['Summarization', '4.2%'],
    #               ['Classification', '3.5%'],
    #               ['Other', '3.5%'],
    #               ['Closed QA', '2.6%'],
    #               ['Extract', '1.9%']]
    # pdb.set_trace()
    replace_image_with_table(slide_page, shape_idx, table_data)
    # pdb.set_trace()
    merge_cells(slide_page, 0, [(0,0), (1,0)])
    merge_cells(slide_page, 0, [(0,1), (1,1)])
    # merge_cell(slide_page, 0, [(1,0), (1,1)])
    print(slide_page.shapes)
    print(1)
    presentation.save('modified.pptx')
    # pdb.set_trace()
    print("Image replaced with table and saved to modified.pptx")